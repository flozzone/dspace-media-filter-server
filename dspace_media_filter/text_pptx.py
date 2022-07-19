import pptx.exc

from dspace_media_filter.filter import MediaFilterRequest, MediaFilterResponse
from pptx import Presentation

from dspace_media_filter.text import TextFilter


class FilterModule(TextFilter):
    def filter(self, req: MediaFilterRequest) -> MediaFilterResponse:
        try:
            return super(FilterModule, self).filter(req)
        except pptx.exc.PackageNotFoundError as e:
            return MediaFilterResponse(error=f"PPTX extraction could not detect the file as a valid PPTX file")
        except pptx.exc.PythonPptxError as e:
            return MediaFilterResponse(error=f"PPTX extraction failed with error {e.__class__.__name__}")

    def filter_text(self, req: MediaFilterRequest) -> str:
        # the following could be also rewritten as a "one"-liner
        # return '\n'.join(['\n'.join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
        #                             for slide in Presentation(req.abs_file).slides])

        text = ""
        prs = Presentation(req.abs_file)

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text

        return text
